"""Tier A text extraction paths and helpers (PDF, rendered HTML/PDF, office)."""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from typing import Any

import pymupdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from scraper.extract_pdf_text import text_path_for_pdf

OCR_DPI = 200
MIN_CHARS_EMPTY = 40


def text_output_path(row: dict[str, Any], output_dir: Path) -> Path:
    """Target text file path for a library row."""
    local_text = row.get("local_text_path")
    if local_text:
        return Path(local_text)

    local_path = row.get("local_path")
    if not local_path:
        raise ValueError(f"Row missing local_path: {row.get('source_url')}")

    asset_path = Path(local_path)
    kind = row.get("library_kind")

    if kind == "source_pdf":
        return text_path_for_pdf(asset_path, output_dir)
    if kind == "rendered_pdf":
        parts = asset_path.parts
        idx = parts.index("rendered")
        return output_dir / "text" / Path(*parts[idx + 1 :]).with_suffix(".txt")
    if kind == "office":
        parts = asset_path.parts
        if "office" in parts:
            idx = parts.index("office")
            relative = Path(*parts[idx + 1 :])
        else:
            relative = Path(asset_path.name)
        return output_dir / "text" / "office" / relative.with_suffix(".txt")
    raise ValueError(f"Unsupported library_kind: {kind}")


def extract_plain_pdf(pdf_path: Path) -> str:
    parts: list[str] = []
    with pymupdf.open(pdf_path) as doc:
        for page in doc:
            text = page.get_text("text")
            if text and text.strip():
                parts.append(text.strip())
    return "\n\n".join(parts).strip()


def extract_pdf_with_ocr(pdf_path: Path, *, dpi: int = OCR_DPI) -> str:
    parts: list[str] = []
    with pymupdf.open(pdf_path) as doc:
        for page in doc:
            text_page = page.get_textpage_ocr(dpi=dpi, full=True)
            text = page.get_text(textpage=text_page) or ""
            if text.strip():
                parts.append(text.strip())
    return "\n\n".join(parts).strip()


def extract_pdf_text(pdf_path: Path, *, use_ocr: bool = False) -> str:
    text = extract_plain_pdf(pdf_path)
    if len(text) >= MIN_CHARS_EMPTY or not use_ocr:
        return text
    ocr_text = extract_pdf_with_ocr(pdf_path)
    return ocr_text if len(ocr_text) > len(text) else text


def extract_docx(path: Path) -> str:
    document = DocxDocument(str(path))
    return "\n\n".join(p.text.strip() for p in document.paragraphs if p.text.strip()).strip()


def extract_xlsx(path: Path) -> str:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"## {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(parts).strip()


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            parts.append(f"## Slide {slide_index}\n" + "\n".join(lines))
    return "\n\n".join(parts).strip()


def extract_legacy_doc(path: Path) -> str:
    if shutil.which("textutil"):
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True,
            check=False,
        )
        if result.returncode == 0:
            return result.stdout.decode("utf-8", errors="replace").strip()
    return ""


def extract_office(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".doc":
        return extract_legacy_doc(path)
    return ""


def extract_for_row(row: dict[str, Any], output_dir: Path) -> tuple[str, str, Path]:
    """Return (text, extraction_method, output_path)."""
    out_path = text_output_path(row, output_dir)
    kind = row.get("library_kind")
    local_path = Path(row["local_path"])

    if kind in {"source_pdf", "rendered_pdf"}:
        existing = out_path if out_path.exists() else None
        if existing and len(existing.read_text(encoding="utf-8", errors="replace").split()) >= 50:
            text = existing.read_text(encoding="utf-8", errors="replace")
            return text, "existing_html_or_text", out_path

        text = extract_pdf_text(local_path, use_ocr=True)
        method = "pdf_text"
        if len(text) < MIN_CHARS_EMPTY:
            method = "pdf_ocr" if text else "pdf_empty"
        return text, method, out_path

    if kind == "office":
        text = extract_office(local_path)
        suffix = local_path.suffix.lower()
        method = f"office{suffix}"
        return text, method, out_path

    raise ValueError(f"Unsupported library_kind: {kind}")
